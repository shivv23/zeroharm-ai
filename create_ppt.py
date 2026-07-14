from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK = RGBColor(0x0A, 0x12, 0x21)
ACCENT = RGBColor(0x00, 0x9E, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xCC, 0xDD, 0xEE)
GRAY = RGBColor(0x88, 0x99, 0xAA)
GREEN = RGBColor(0x00, 0xC8, 0x53)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
RED = RGBColor(0xFF, 0x33, 0x33)

def bg(slide, color=DARK):
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = color
    bg_shape.line.fill.background()
    slide.shapes._spTree.remove(bg_shape._element)
    slide.shapes._spTree.insert(2, bg_shape._element)
    return bg_shape

def accent_line(slide, left, top, width, color=ACCENT, thickness=None):
    h = thickness if thickness is not None else Pt(4)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_frame(slide, left, top, width, height, items, font_size=16, color=LIGHT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(8)
        p.level = 0
    return txBox

def add_card(slide, left, top, width, height, title, body, title_color=ACCENT):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x0F, 0x1A, 0x2E)
    card.line.color.rgb = RGBColor(0x1A, 0x2A, 0x44)
    card.line.width = Pt(1)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4), title, font_size=14, bold=True, color=title_color)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6), body, font_size=11, color=LIGHT)
    return card

def slide_number(slide, num, total):
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4), f"{num}/{total}", font_size=10, color=GRAY, alignment=PP_ALIGN.RIGHT)

TOTAL = 13

# ==============================
# SLIDE 1: TITLE
# ==============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
accent_line(slide, Inches(1), Inches(2.6), Inches(3), ACCENT)
add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(1.2), "ZeroHarm AI", font_size=54, bold=True, color=WHITE)
add_text_box(slide, Inches(1), Inches(3.9), Inches(11), Inches(0.6), "Industrial Safety Intelligence Platform", font_size=26, color=LIGHT)
add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5), "ET AI Hackathon 2026 — Problem Statement 1", font_size=16, color=GRAY)
accent_line(slide, Inches(1), Inches(5.6), Inches(11), color=RGBColor(0x1A, 0x2A, 0x44), thickness=Pt(1))
add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.4), "Created by Shivam Kumar", font_size=18, bold=True, color=ACCENT)
slide_number(slide, 1, TOTAL)

# ==============================
# SLIDE 2: THE PROBLEM
# ==============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
accent_line(slide, Inches(0.8), Inches(0.5), Inches(2), ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.6), "The Problem", font_size=36, bold=True, color=WHITE)

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.0),
    "On January 20, 2025, the Visakhapatnam Steel Plant suffered a catastrophic explosion. "
    "Eight workers died. Twelve were injured. The root cause — gas pressure sensors showed "
    "warning signals, but no intelligence layer connected those readings to operational decisions.",
    font_size=16, color=LIGHT)

add_card(slide, Inches(0.8), Inches(3.0), Inches(3.6), Inches(2.0),
    "Siloed Data", "Sensor data, permit systems, incident reports, and compliance logs exist in isolation. No unified view of plant safety.")
add_card(slide, Inches(4.8), Inches(3.0), Inches(3.6), Inches(2.0),
    "Reactive Response", "Incidents are caught after they happen. Root cause analysis is manual and slow. No predictive capability.")
add_card(slide, Inches(8.8), Inches(3.0), Inches(3.6), Inches(2.0),
    "No Safe Simulation", "No way to model 'what-if' scenarios without risking real operations. Emergency drills are disconnected from live data.")

add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.5),
    '"Sensor data without an intelligence layer is just noise."',
    font_size=18, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
slide_number(slide, 2, TOTAL)

# ==============================
# SLIDE 3: THE SOLUTION
# ==============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
accent_line(slide, Inches(0.8), Inches(0.5), Inches(2), ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.6), "The Solution", font_size=36, bold=True, color=WHITE)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.2),
    "ZeroHarm AI bridges the gap between raw industrial sensor data and safety-critical decisions through a multi-agent architecture "
    "that detects compound risks — where multiple minor issues combine into a critical threat — before they become catastrophes.",
    font_size=16, color=LIGHT)

bullets = [
    "Compound Risk Detection Engine — correlates sensor readings, permit activity, and maintenance status in real time",
    "Digital Twin Aggregator — unified plant health dashboard with 25+ live KPI metrics",
    "Real-time WebSocket Push — 2-second update cycle to all connected clients",
    "Machine Learning — Isolation Forest anomaly detection & predictive risk forecasting with graceful fallback",
    "Interactive Simulation — What-If scenario modeling and Emergency Response orchestration",
    "Regulatory Compliance — Automated audit, reporting, and CAPA workflows against OSHA, EPA, ISO, OISD standards",
]
add_bullet_frame(slide, Inches(0.8), Inches(2.9), Inches(11.5), Inches(4.0), bullets, font_size=15, color=LIGHT)
slide_number(slide, 3, TOTAL)

# ==============================
# SLIDE 4: ARCHITECTURE
# ==============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
accent_line(slide, Inches(0.8), Inches(0.5), Inches(2), ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.6), "System Architecture", font_size=36, bold=True, color=WHITE)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(3.5), Inches(0.4), "Presentation Layer", font_size=16, bold=True, color=ACCENT)
add_text_box(slide, Inches(0.8), Inches(1.9), Inches(3.5), Inches(1.0), "React 18 + WebSocket\nGeospatial SVG Heatmap, Dashboards, Chat Widget\nReal-time state via WebSocket push", font_size=12, color=LIGHT)
shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.5), Inches(1.8), Inches(0.6), Inches(0.3))
shape.fill.solid(); shape.fill.fore_color.rgb = GRAY; shape.line.fill.background()

add_text_box(slide, Inches(5.3), Inches(1.5), Inches(3.5), Inches(0.4), "API Gateway", font_size=16, bold=True, color=ACCENT)
add_text_box(slide, Inches(5.3), Inches(1.9), Inches(3.5), Inches(1.0), "FastAPI + Uvicorn\n55+ REST Endpoints, WebSocket Server\nJWT Auth, RBAC, Rate Limiting", font_size=12, color=LIGHT)
shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(9.0), Inches(1.8), Inches(0.6), Inches(0.3))
shape.fill.solid(); shape.fill.fore_color.rgb = GRAY; shape.line.fill.background()

add_text_box(slide, Inches(9.8), Inches(1.5), Inches(3.0), Inches(0.4), "Orchestration Layer", font_size=16, bold=True, color=ACCENT)
add_text_box(slide, Inches(9.8), Inches(1.9), Inches(3.0), Inches(1.0), "Multi-Agent Python Engine\nCompound Risk, Emergency Response\nWhat-If Simulator, 15+ Agents", font_size=12, color=LIGHT)

agents_left = [
    "Sensor Monitor Agent", "Permit Activity Agent", "Maintenance Status Agent",
    "Fusion Supervisor", "Quality Compliance Agent", "Anomaly Detector",
    "Predictive Risk Forecaster", "Root Cause Analyzer", "Digital Twin Aggregator",
]
agents_right = [
    "Emergency Response Orchestrator", "Incident Pattern Intelligence", "Alert Triage Engine",
    "Personnel Tracker", "Environmental Monitor", "Equipment Health Monitor",
    "Safety Observation System", "Chat Assistant", "Report Generator (ReportLab)",
]
add_text_box(slide, Inches(0.8), Inches(3.3), Inches(11.5), Inches(0.3), "18 Agent Modules", font_size=14, bold=True, color=ACCENT)
add_bullet_frame(slide, Inches(0.8), Inches(3.7), Inches(5.5), Inches(3.2), [f"\u2022  {a}" for a in agents_left], font_size=13, color=LIGHT)
add_bullet_frame(slide, Inches(6.8), Inches(3.7), Inches(5.5), Inches(3.2), [f"\u2022  {a}" for a in agents_right], font_size=13, color=LIGHT)
slide_number(slide, 4, TOTAL)

# ==============================
# SLIDE 5: FEATURES — CORE
# ==============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
accent_line(slide, Inches(0.8), Inches(0.5), Inches(2), ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.6), "Key Features — Core Modules", font_size=36, bold=True, color=WHITE)

features1 = [
    ("Compound Risk Detection", "Multi-factor correlation of sensor, permit, and maintenance data to detect cascading failure risks"),
    ("Geospatial Risk Heatmap", "10-zone SVG plant layout with color-coded risk overlay, sensor dots, and permit badges"),
    ("Emergency Response Orchestration", "Full trigger, mustering, event log, and resolution lifecycle with real-time timeline"),
    ("What-If Simulator", "Pre-built and custom scenario modeling with live risk impact — gas leak, fire, power loss, etc."),
    ("Digital Twin Dashboard", "Unified plant pulse with KPI gauges — health, risk, sensors, compliance, zones, alerts"),
    ("Root Cause Analysis", "Causal chain engine with 5-Why analysis, confidence scoring, and prioritized recommendations"),
]
for i, (title, desc) in enumerate(features1):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.6 + row * 2.6)
    add_card(slide, x, y, Inches(3.8), Inches(2.2), title, desc)

slide_number(slide, 5, TOTAL)

# ==============================
# SLIDE 6: FEATURES — ADVANCED
# ==============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
accent_line(slide, Inches(0.8), Inches(0.5), Inches(2), ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.6), "Key Features — Advanced Modules", font_size=36, bold=True, color=WHITE)

features2 = [
    ("Cost of Safety Dashboard", "Financial translation: total cost, regulatory fines, risk exposure, prevention savings with charts"),
    ("Anomaly Detection", "Isolation Forest on sensor history with z-score confidence — flags unusual patterns in real time"),
    ("Predictive Risk Forecast", "5-step trend projection with direction indicators (rising / stable / falling)"),
    ("Safety Gamification", "Per-zone leaderboard, trophies for top zones, warning flags for bottom zones"),
    ("Alert Triage Engine", "Per-alert urgency scoring with prioritized response actions — evacuate / dispatch / isolate"),
    ("AI Chat Assistant", "Natural-language query interface with live plant context injection"),
    ("Personnel Tracker", "Zone occupancy, hazard exposure tracking, emergency mustering drills with progress"),
    ("Environmental Monitoring", "Per-zone air quality — CO, NOx, SO2, PM, VOCs — against regulatory limits"),
    ("Equipment Health Monitor", "Real-time health scores for critical assets with remaining useful life estimates"),
    ("Safety Observations", "Worker-submitted hazard reports with type/severity classification and trend analytics"),
    ("Compliance Audit Agent", "Automated 32-check audits with weighted scoring and regulatory standard mapping"),
    ("Regulatory Reporter", "Automated PDF report generation against OSHA, EPA, ISO, NFPA, OISD standards"),
]
for i, (title, desc) in enumerate(features2):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.4 + row * 1.4)
    txBox = slide.shapes.add_textbox(x, y, Inches(3.9), Inches(1.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.name = "Calibri"
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(9)
    p2.font.color.rgb = LIGHT
    p2.font.name = "Calibri"
    p2.space_before = Pt(2)

slide_number(slide, 6, TOTAL)

# ==============================
# SLIDE 7: TECH STACK
# ==============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
accent_line(slide, Inches(0.8), Inches(0.5), Inches(2), ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.6), "Technology Stack", font_size=36, bold=True, color=WHITE)

stack = [
    ("Frontend", "React 18, Webpack 5, SVG, CSS-in-JS"),
    ("Backend", "Python 3.13, FastAPI, Uvicorn"),
    ("Communication", "WebSocket (real-time push, 2s interval) + REST API"),
    ("Machine Learning", "scikit-learn (Isolation Forest, Linear Regression — graceful fallback)"),
    ("Database", "SQLite + SQLAlchemy async (aiosqlite)"),
    ("Authentication", "JWT (PyJWT) + bcrypt password hashing"),
    ("PDF Generation", "ReportLab"),
    ("Containerization", "Docker, docker-compose"),
    ("Deployment", "AWS EC2 (t2.micro, Docker)"),
    ("Data Generation", "NumPy-based synthetic simulation engine"),
]
for i, (layer, tech) in enumerate(stack):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.5 + row * 1.1)
    add_text_box(slide, x, y, Inches(1.8), Inches(0.4), layer, font_size=13, bold=True, color=ACCENT)
    add_text_box(slide, x + Inches(1.9), y, Inches(4.0), Inches(0.4), tech, font_size=13, color=LIGHT)
    if i < len(stack) - 1:
        accent_line(slide, x, y + Inches(0.45), Inches(5.8), color=RGBColor(0x1A, 0x2A, 0x44), thickness=Pt(1))

slide_number(slide, 7, TOTAL)

# ==============================
# SLIDE 8: DATA FLOW
# ==============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
accent_line(slide, Inches(0.8), Inches(0.5), Inches(2), ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.6), "Data Flow & Simulation Loop", font_size=36, bold=True, color=WHITE)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
    "The entire platform runs on a continuous 2-second simulation loop:", font_size=16, color=LIGHT)

steps = [
    "1. SyntheticDataGenerator.step() produces fresh sensor readings, permit states, zone risks",
    "2. Plant state persisted to SQLite via save_plant_state, save_sensor_reading, save_permit",
    "3. CompoundRiskDetectionEngine.run_async() fuses sensor + permit + maintenance data",
    "4. Risk trend updated, anomaly detector fed, predictive forecaster updated",
    "5. Agent Activity Feed logs each step — sensor scan, permit audit, maintenance check",
    "6. WebSocket broadcasts full state + risk to all connected frontend clients",
    "7. Alert dispatcher sends notifications for high-severity events (console / Slack / email / SMS)",
]
add_bullet_frame(slide, Inches(0.8), Inches(2.2), Inches(11.5), Inches(3.5), steps, font_size=14, color=LIGHT)

add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.8),
    "The WebSocket reconnects with exponential backoff (1s base, 30s max, 20 retries). "
    "A 3-second debounce prevents brief reconnection flashes on the disconnected banner.",
    font_size=13, color=GRAY)
slide_number(slide, 8, TOTAL)

# ==============================
# SLIDE 9: AUTH & SECURITY
# ==============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
accent_line(slide, Inches(0.8), Inches(0.5), Inches(2), ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.6), "Authentication & Security", font_size=36, bold=True, color=WHITE)

add_card(slide, Inches(0.8), Inches(1.5), Inches(3.8), Inches(2.5),
    "JWT Authentication",
    "Token-based auth with HS256 signing\n"
    "Access tokens: configurable expiry (default 24h)\n"
    "Refresh tokens: 7-day rotation, revocable\n"
    "Automatic injection via authFetch.js interceptor")

add_card(slide, Inches(4.8), Inches(1.5), Inches(3.8), Inches(2.5),
    "Role-Based Access Control",
    "4 roles with granular permissions:\n"
    "  admin — full access + user management\n"
    "  safety_officer — read/write/emergency\n"
    "  operator — read/write permits\n"
    "  viewer — read-only")

add_card(slide, Inches(8.8), Inches(1.5), Inches(3.8), Inches(2.5),
    "Security Hardening",
    "API key verification on write endpoints\n"
    "Rate limiting (100 req/min per IP)\n"
    "Audit trail for all auth events\n"
    "No secrets in code — all via .env\n"
    "Default JWT_SECRET warning at startup")

add_card(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(2.0),
    "Default Credentials (for evaluation)",
    "admin / admin123 — full admin access\n"
    "safety_officer / safety_officer123 — safety operations\n"
    "operator / operator123 — shift operator\n"
    "viewer / viewer123 — read-only")

add_card(slide, Inches(6.8), Inches(4.5), Inches(5.8), Inches(2.0),
    "WebSocket Security",
    "Token passed as ?token= query parameter\n"
    "Server validates JWT on every connection\n"
    "Exponential backoff reconnection (20 retries)\n"
    "Intentional disconnect flag prevents reconnect storms")

slide_number(slide, 9, TOTAL)

# ==============================
# SLIDE 10: KEY RESOLVED ISSUES
# ==============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
accent_line(slide, Inches(0.8), Inches(0.5), Inches(2), ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.6), "Key Technical Challenges Solved", font_size=36, bold=True, color=WHITE)

issues = [
    ("Chunk Loading Failures", "Removed React.lazy() — all components imported directly to avoid webpack chunk load errors"),
    ("API Authorization (401)", "Global authFetch.js interceptor auto-injects Bearer token on every /api/ request"),
    ("CORS in Development", "Added 'proxy': 'http://localhost:8000' to frontend package.json for seamless dev"),
    ("Backend 500 — Digital Twin", "Fixed _trend_direction: was multiplying dict by float — now extracts score key correctly"),
    ("Anomalies Endpoint Hanging", "Wrapped IsolationForest scan_all() in asyncio.to_thread() to unblock event loop"),
    ("Emergency Log 'undefined'", "entry.msg -> entry.text — corrected response field name"),
    ("Dashboard Loading Flash", "initialLoad ref prevents 'Loading...' state on 5-second interval refetch"),
    ("Disconnected Banner Flash", "3-second debounce before showing 'Disconnected — showing cached data'"),
    ("Data Binding — Safety Scores", "data.data.zones -> data.zones — API returns zones at root, not nested"),
]
for i, (title, desc) in enumerate(issues):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.5 + row * 1.2)
    add_text_box(slide, x, y, Inches(6.0), Inches(0.3), f"\u2022  {title}", font_size=13, bold=True, color=ORANGE)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.3), Inches(5.8), Inches(0.7), desc, font_size=11, color=LIGHT)

slide_number(slide, 10, TOTAL)

# ==============================
# SLIDE 11: API ENDPOINTS
# ==============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
accent_line(slide, Inches(0.8), Inches(0.5), Inches(2), ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.6), "API Surface (55+ Endpoints)", font_size=36, bold=True, color=WHITE)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(3.5), Inches(0.4), "Interactive Swagger docs at /docs", font_size=14, bold=True, color=ACCENT)

categories = [
    ("Health & Auth", "/api/health, /api/auth/login, /api/auth/register, /api/auth/me, /api/auth/refresh"),
    ("Plant & Sensors", "/api/plant/state, /api/sensors, /api/sensors/zone/{id}"),
    ("Risk & Alerts", "/api/risk/current, /api/risk/alerts, /api/risk-trend, /api/alerts/triage"),
    ("Permits & Compliance", "/api/permits, /api/compliance/audit, /api/compliance/trend, /api/health-index"),
    ("Digital Twin & Analytics", "/api/digital-twin, /api/cost-of-safety, /api/anomalies, /api/forecast"),
    ("Emergency & What-If", "/api/emergency/trigger, /api/emergency/active, /api/what-if/scenarios, /api/what-if/apply"),
    ("Personnel & Environment", "/api/personnel/locations, /api/environmental/metrics, /api/safety/observations"),
    ("Admin & Reports", "/api/admin/users, /api/audit/log, /api/reports/regulatory/{std}, /api/capa/statistics"),
    ("Chat & Knowledge", "/api/chat, /api/kg/query, /api/activity-feed, /api/incident-patterns"),
]
for i, (cat, eps) in enumerate(categories):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(2.0 + row * 1.7)
    add_text_box(slide, x, y, Inches(3.9), Inches(0.3), cat, font_size=12, bold=True, color=ACCENT)
    add_text_box(slide, x, y + Inches(0.3), Inches(3.9), Inches(1.2), eps, font_size=10, color=LIGHT)

add_text_box(slide, Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.4),
    "All endpoints protected by JWT + RBAC. Write endpoints additionally require X-API-Key header. Rate limited to 100 req/min/IP.",
    font_size=11, color=GRAY)
slide_number(slide, 11, TOTAL)

# ==============================
# SLIDE 12: TESTING & DEPLOYMENT
# ==============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
accent_line(slide, Inches(0.8), Inches(0.5), Inches(2), ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.6), "Testing & Deployment", font_size=36, bold=True, color=WHITE)

add_card(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5),
    "Test Suite (56 pytest tests)",
    "Compound Risk Engine — risk scoring, alert generation, severity levels\n"
    "Quality Compliance — 32-check audit, severity mapping, weighted scoring\n"
    "What-If Simulator — scenario application, reset, custom overrides\n"
    "Emergency Response — trigger, mustering, resolution lifecycle\n"
    "Incident Patterns — type/zone/cause discovery, trend statistics\n"
    "Digital Twin — KPI aggregation, trend direction analysis\n"
    "Anomaly Detection — Isolation Forest scoring, z-score computation\n"
    "Root Cause Analysis — causal chain, confidence scoring\n"
    "Synthetic Data Generator — sensor values, permit lifecycle, zone risks")

add_card(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.5),
    "Deployment Options",
    "Docker: docker compose up -d (backend :8000, frontend :8080)\n"
    "Manual: pip install + uvicorn (backend), npm install + npm start (frontend)\n"
    "AWS: http://13.59.83.134 (EC2 t2.micro, Docker Compose)\n\n"
    "Requirements:\n"
    "  Python 3.13+, Node 20+\n"
    "  Core: fastapi, uvicorn, sqlalchemy, aiosqlite, numpy, pyjwt, bcrypt\n"
    "  ML (optional): scikit-learn for anomaly detection & forecasting\n"
    "  Vision (optional): opencv-python, ultralytics for camera integration")

add_card(slide, Inches(0.8), Inches(4.5), Inches(11.8), Inches(2.5),
    "Project Structure",
    "backend/ — 30+ Python modules (agents, orchestrators, ML models, auth, database)\n"
    "frontend/ — 25+ React components, WebSocket store, auth interceptor, theme system\n"
    "Root: docker-compose.yml, render.yaml, Makefile",
    title_color=GREEN)

slide_number(slide, 12, TOTAL)

# ==============================
# SLIDE 13: THANK YOU
# ==============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
accent_line(slide, Inches(4.5), Inches(2.5), Inches(4.5), ACCENT)
add_text_box(slide, Inches(1), Inches(2.8), Inches(11.3), Inches(1.2), "Thank You", font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.9), Inches(11.3), Inches(0.6), "ZeroHarm AI — Industrial Safety Intelligence Platform", font_size=22, color=LIGHT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.8), Inches(11.3), Inches(0.5), "ET AI Hackathon 2026", font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.5), Inches(11.3), Inches(0.5), "Created by Shivam Kumar", font_size=20, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(6.2), Inches(11.3), Inches(0.4), "Live Demo: http://13.59.83.134", font_size=14, color=ACCENT, alignment=PP_ALIGN.CENTER)
slide_number(slide, 13, TOTAL)

output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ZeroHarm_AI_Presentation.pptx")
prs.save(output_path)
print(f"Saved to: {output_path}")
